import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def apply_text_styling(paragraph, font_name="Arial", size_pt=18, color_rgb=(51, 65, 85), bold=False, italic=False):
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size_pt)
    paragraph.font.color.rgb = RGBColor(*color_rgb)
    paragraph.font.bold = bold
    paragraph.font.italic = italic

def main():
    if len(sys.argv) < 3:
        print("Usage: python3 ppt_exporter.py <slides_json_or_file> <output_pptx_path>")
        sys.exit(1)

    json_input = sys.argv[1]
    output_path = sys.argv[2]

    # Load JSON from file or parse directly
    try:
        if os.path.exists(json_input):
            with open(json_input, 'r', encoding='utf-8') as f:
                slides = json.load(f)
        else:
            slides = json.loads(json_input)
    except Exception as e:
        print(f"Error loading JSON input: {e}", file=sys.stderr)
        sys.exit(1)

    # Allow slides key wrapping or direct array
    if isinstance(slides, dict) and "slides" in slides:
        slides = slides["slides"]

    prs = Presentation()
    # Set slide size to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout for custom placement to make it look premium
    blank_layout = prs.slide_layouts[6]

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        
        title_text = slide_data.get("title", f"Slide {idx + 1}")
        content_points = slide_data.get("content", [])
        visual_suggestion = slide_data.get("visual_suggestion", "")

        if idx == 0:
            # Title Slide: Dark Slate theme (#0F172A)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(15, 23, 42)
            
            # Text box for title & description
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p_title = tf.paragraphs[0]
            p_title.text = title_text
            apply_text_styling(p_title, font_name="Georgia", size_pt=44, color_rgb=(255, 255, 255), bold=True)
            p_title.space_after = Pt(20)
            
            if content_points:
                p_sub = tf.add_paragraph()
                p_sub.text = " • ".join(content_points)
                apply_text_styling(p_sub, font_name="Arial", size_pt=18, color_rgb=(148, 163, 184))
        else:
            # Content Slide: Light slate background (#F8FAFC)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(248, 250, 252)
            
            # Slide Header Title
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            apply_text_styling(p_title, font_name="Georgia", size_pt=32, color_rgb=(15, 23, 42), bold=True)
            
            # Left Column: Content Points
            content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.8))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            for p_idx, point in enumerate(content_points):
                p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                p.text = f"•  {point}"
                apply_text_styling(p, font_name="Arial", size_pt=16, color_rgb=(51, 65, 85))
                p.space_after = Pt(14)
                
            # Right Column: Visual Layout Recommendation Box
            if visual_suggestion:
                visual_box = slide.shapes.add_textbox(Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.8))
                tf_visual = visual_box.text_frame
                tf_visual.word_wrap = True
                
                p_vis_lbl = tf_visual.paragraphs[0]
                p_vis_lbl.text = "RECOMMENDED VISUAL LAYOUT:"
                apply_text_styling(p_vis_lbl, font_name="Arial", size_pt=12, color_rgb=(99, 102, 241), bold=True)
                p_vis_lbl.space_after = Pt(10)
                
                p_vis_val = tf_visual.add_paragraph()
                p_vis_val.text = visual_suggestion
                apply_text_styling(p_vis_val, font_name="Arial", size_pt=14, color_rgb=(71, 85, 105), italic=True)

    prs.save(output_path)
    print(f"PowerPoint saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
