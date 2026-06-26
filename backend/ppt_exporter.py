import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def apply_text_styling(paragraph, font_name="Arial", size_pt=18, color_rgb=(51, 65, 85), bold=False, italic=False):
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size_pt)
    paragraph.font.color.rgb = RGBColor(*color_rgb)
    paragraph.font.bold = bold
    paragraph.font.italic = italic

def add_image_from_url(slide, url, left, top, width, height):
    if not url:
        return
    import urllib.request
    import tempfile
    import os
    import ssl
    
    try:
        # Bypassing SSL verification on macOS
        ssl_context = ssl._create_unverified_context()
        # Add User-Agent to prevent 403 blocks
        req = urllib.request.Request(
            url,
            headers={
                'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
            }
        )
        # Download image to a temp file
        with urllib.request.urlopen(req, context=ssl_context, timeout=15) as response:
            with tempfile.NamedTemporaryFile(delete=False, suffix=".jpg") as temp_file:
                temp_file.write(response.read())
                temp_file_path = temp_file.name
        
        # Add to slide
        slide.shapes.add_picture(temp_file_path, left, top, width=width, height=height)
        
        # Clean up temp file
        os.unlink(temp_file_path)
    except Exception as e:
        print(f"Error adding image from {url}: {e}", file=sys.stderr)

THEME_PALETTES = {
    "slate": {
        "bg": (248, 250, 252),
        "cardBg": (255, 255, 255),
        "text": (15, 23, 42),
        "muted": (100, 116, 139),
        "accent": (79, 70, 229),
        "accentLight": (239, 246, 255),
        "border": (226, 232, 240)
    },
    "dark": {
        "bg": (15, 23, 42),
        "cardBg": (30, 41, 59),
        "text": (248, 250, 252),
        "muted": (148, 163, 184),
        "accent": (129, 140, 248),
        "accentLight": (49, 46, 129),
        "border": (51, 65, 85)
    },
    "indigo": {
        "bg": (238, 242, 255),
        "cardBg": (255, 255, 255),
        "text": (30, 27, 75),
        "muted": (79, 70, 229),
        "accent": (79, 70, 229),
        "accentLight": (224, 231, 255),
        "border": (199, 210, 254)
    },
    "emerald": {
        "bg": (240, 253, 244),
        "cardBg": (255, 255, 255),
        "text": (6, 78, 59),
        "muted": (5, 150, 105),
        "accent": (5, 150, 105),
        "accentLight": (209, 250, 229),
        "border": (167, 243, 208)
    },
    "sakura": {
        "bg": (255, 245, 245),
        "cardBg": (255, 255, 255),
        "text": (76, 5, 25),
        "muted": (219, 39, 119),
        "accent": (219, 39, 119),
        "accentLight": (252, 231, 243),
        "border": (251, 207, 232)
    },
    "cobalt": {
        "bg": (240, 244, 248),
        "cardBg": (255, 255, 255),
        "text": (16, 42, 67),
        "muted": (15, 98, 254),
        "accent": (15, 98, 254),
        "accentLight": (229, 241, 255),
        "border": (208, 226, 255)
    }
}

def draw_diagram(slide, suggestion, content, diagram=None, theme_name="slate", x_shift=Inches(0)):
    text = (suggestion or "").lower()
    
    # Helper to truncate text inside shapes
    def trunc(s, max_len=30):
        return s[:max_len] + "..." if len(s) > max_len else s

    palette = THEME_PALETTES.get(theme_name, THEME_PALETTES["slate"])
    accent_rgb = palette["accent"]
    accent_light_rgb = palette["accentLight"]
    text_dark_rgb = palette["text"]
    text_light_rgb = (255, 255, 255)
    border_rgb = palette["border"]

    # Draw structured diagram if available
    if diagram and isinstance(diagram, dict) and diagram.get("type") and diagram.get("type") != "none":
        diagram_type = diagram["type"]
        data = diagram.get("data", {})
        items = data.get("items", [])
        
        # 1. Flowchart
        if diagram_type == "flowchart" and items:
            top_offset = Inches(1.8)
            for i, item in enumerate(items[:4]):
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3) + x_shift, top_offset, Inches(3.8), Inches(0.9))
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
                shape.line.color.rgb = RGBColor(*border_rgb)
                tf = shape.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get("label", f"Step {i+1}")
                apply_text_styling(p, font_name="Arial", size_pt=11, color_rgb=accent_rgb, bold=True)
                
                if item.get("detail"):
                    p2 = tf.add_paragraph()
                    p2.text = trunc(item["detail"], 70)
                    apply_text_styling(p2, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
                    
                top_offset += Inches(0.9)
                
                if i < len(items[:4]) - 1:
                    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.1) + x_shift, top_offset, Inches(0.2), Inches(0.2))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(*accent_rgb)
                    arrow.line.fill.background()
                    top_offset += Inches(0.2)
            return

        # 2. Timeline
        elif diagram_type == "timeline" and items:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5) + x_shift, Inches(1.8), Inches(0.08), Inches(4.3))
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(*accent_rgb)
            bar.line.fill.background()
            
            top_offset = Inches(1.9)
            for i, item in enumerate(items[:4]):
                circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.34) + x_shift, top_offset + Inches(0.12), Inches(0.4), Inches(0.4))
                circle.fill.solid()
                circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
                circle.line.color.rgb = RGBColor(*accent_rgb)
                circle.line.width = Pt(3)
                
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9) + x_shift, top_offset, Inches(3.2), Inches(0.85))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
                card.line.color.rgb = RGBColor(*border_rgb)
                tf = card.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get("label", f"Phase {i+1}")
                apply_text_styling(p, font_name="Arial", size_pt=10, color_rgb=accent_rgb, bold=True)
                
                if item.get("detail"):
                    p2 = tf.add_paragraph()
                    p2.text = trunc(item["detail"], 70)
                    apply_text_styling(p2, font_name="Arial", size_pt=8.5, color_rgb=text_dark_rgb)
                    
                top_offset += Inches(1.1)
            return

        # 3. Organizational Chart / Hierarchy
        elif diagram_type == "org_chart":
            root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2) + x_shift, Inches(1.8), Inches(2.0), Inches(0.7))
            root.fill.solid()
            root.fill.fore_color.rgb = RGBColor(*accent_rgb)
            root.line.fill.background()
            p_root = root.text_frame.paragraphs[0]
            p_root.text = data.get("root", "Root Level")
            apply_text_styling(p_root, font_name="Arial", size_pt=11, color_rgb=text_light_rgb, bold=True)
            
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.17) + x_shift, Inches(2.5), Inches(0.05), Inches(0.5))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*accent_rgb)
            line.line.fill.background()
            
            hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3) + x_shift, Inches(3.0), Inches(3.8), Inches(0.05))
            hbar.fill.solid()
            hbar.fill.fore_color.rgb = RGBColor(*accent_rgb)
            hbar.line.fill.background()
            
            for i, item in enumerate(items[:3]):
                x_pos = Inches(7.8) + x_shift + i * Inches(1.4)
                vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.57), Inches(3.05), Inches(0.05), Inches(0.35))
                vline.fill.solid()
                vline.fill.fore_color.rgb = RGBColor(*accent_rgb)
                vline.line.fill.background()
                
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.4), Inches(1.2), Inches(1.4))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
                card.line.color.rgb = RGBColor(*border_rgb)
                tf = card.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get("label", f"Div {i+1}")
                apply_text_styling(p, font_name="Arial", size_pt=9, color_rgb=accent_rgb, bold=True)
                
                if item.get("detail"):
                    p2 = tf.add_paragraph()
                    p2.text = trunc(item["detail"], 40)
                    apply_text_styling(p2, font_name="Arial", size_pt=8, color_rgb=text_dark_rgb)
            return

        # 4. Mind Map
        elif diagram_type == "mind_map":
            center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2) + x_shift, Inches(3.3), Inches(1.8), Inches(1.0))
            center.fill.solid()
            center.fill.fore_color.rgb = RGBColor(*accent_rgb)
            center.line.fill.background()
            p_c = center.text_frame.paragraphs[0]
            p_c.text = data.get("root", "Core Focus")
            apply_text_styling(p_c, font_name="Arial", size_pt=11, color_rgb=text_light_rgb, bold=True)
            
            spokes = [
                (Inches(7.8), Inches(1.8)),
                (Inches(10.6), Inches(1.8)),
                (Inches(9.2), Inches(4.8))
            ]
            for i, item in enumerate(items[:3]):
                x, y = spokes[i]
                conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, (x + Inches(9.2))/2 + Inches(0.45) + x_shift, (y + Inches(3.3))/2 + Inches(0.25), Inches(0.05), Inches(0.4))
                conn.fill.solid()
                conn.fill.fore_color.rgb = RGBColor(*border_rgb)
                conn.line.fill.background()
                
                node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + x_shift, y, Inches(1.8), Inches(0.9))
                node.fill.solid()
                node.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
                node.line.color.rgb = RGBColor(*border_rgb)
                tf = node.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get("label", f"Aspect {i+1}")
                apply_text_styling(p, font_name="Arial", size_pt=9, color_rgb=accent_rgb, bold=True)
                
                if item.get("detail"):
                    p2 = tf.add_paragraph()
                    p2.text = trunc(item["detail"], 40)
                    apply_text_styling(p2, font_name="Arial", size_pt=8, color_rgb=text_dark_rgb)
            return

        # 5. Decision Tree
        elif diagram_type == "decision_tree":
            cond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.2) + x_shift, Inches(1.8), Inches(2.0), Inches(1.0))
            cond.fill.solid()
            cond.fill.fore_color.rgb = RGBColor(*accent_rgb)
            cond.line.fill.background()
            p_cond = cond.text_frame.paragraphs[0]
            p_cond.text = data.get("root", "Decision")
            apply_text_styling(p_cond, font_name="Arial", size_pt=10, color_rgb=text_light_rgb, bold=True)
            
            yes_item = next((it for it in items if it.get("category") == "yes"), items[0] if items else {})
            left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8) + x_shift, Inches(3.6), Inches(1.8), Inches(1.4))
            left_card.fill.solid()
            left_card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            left_card.line.color.rgb = RGBColor(34, 197, 94)
            left_card.line.width = Pt(2)
            tf_l = left_card.text_frame
            tf_l.word_wrap = True
            p_l = tf_l.paragraphs[0]
            p_l.text = yes_item.get("label", "YES Path")
            apply_text_styling(p_l, font_name="Arial", size_pt=9, color_rgb=RGBColor(34, 197, 94), bold=True)
            if yes_item.get("detail"):
                p_ld = tf_l.add_paragraph()
                p_ld.text = trunc(yes_item["detail"], 45)
                apply_text_styling(p_ld, font_name="Arial", size_pt=8, color_rgb=text_dark_rgb)
                
            no_item = next((it for it in items if it.get("category") == "no"), items[1] if len(items) > 1 else (items[0] if items else {}))
            right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6) + x_shift, Inches(3.6), Inches(1.8), Inches(1.4))
            right_card.fill.solid()
            right_card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            right_card.line.color.rgb = RGBColor(239, 68, 68)
            right_card.line.width = Pt(2)
            tf_r = right_card.text_frame
            tf_r.word_wrap = True
            p_r = tf_r.paragraphs[0]
            p_r.text = no_item.get("label", "NO Path")
            apply_text_styling(p_r, font_name="Arial", size_pt=9, color_rgb=RGBColor(239, 68, 68), bold=True)
            if no_item.get("detail"):
                p_rd = tf_r.add_paragraph()
                p_rd.text = trunc(no_item["detail"], 45)
                apply_text_styling(p_rd, font_name="Arial", size_pt=8, color_rgb=text_dark_rgb)
            return

        # 6. System Architecture
        elif diagram_type == "architecture" and items:
            top_offset = Inches(1.8)
            for i, item in enumerate(items[:3]):
                bg_rgb = accent_rgb if i % 2 == 1 else accent_light_rgb
                fg_rgb = text_light_rgb if i % 2 == 1 else text_dark_rgb
                lbl_color = text_light_rgb if i % 2 == 1 else accent_rgb
                
                layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3) + x_shift, top_offset, Inches(3.8), Inches(0.85))
                layer.fill.solid()
                layer.fill.fore_color.rgb = RGBColor(*bg_rgb)
                layer.line.color.rgb = RGBColor(*border_rgb)
                tf = layer.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = item.get("label", f"Tier {i+1}")
                apply_text_styling(p, font_name="Arial", size_pt=10, color_rgb=lbl_color, bold=True)
                
                if item.get("detail"):
                    p2 = tf.add_paragraph()
                    p2.text = trunc(item["detail"], 60)
                    apply_text_styling(p2, font_name="Arial", size_pt=8, color_rgb=fg_rgb)
                    
                top_offset += Inches(0.85)
                if i < 2:
                    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.1) + x_shift, top_offset, Inches(0.2), Inches(0.25))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = RGBColor(*accent_rgb)
                    arrow.line.fill.background()
                    top_offset += Inches(0.25)
            return

        # 7. Data Relationship / ERD
        elif diagram_type == "erd" and items:
            item_a = next((it for it in items if it.get("category") == "entity_a"), items[0] if items else {})
            tbl_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8) + x_shift, Inches(2.2), Inches(1.8), Inches(2.4))
            tbl_a.fill.solid()
            tbl_a.fill.fore_color.rgb = RGBColor(255, 255, 255)
            tbl_a.line.color.rgb = RGBColor(*accent_rgb)
            tbl_a.line.width = Pt(2)
            tf_a = tbl_a.text_frame
            tf_a.word_wrap = True
            p_ah = tf_a.paragraphs[0]
            p_ah.text = item_a.get("label", "Entity_A")
            apply_text_styling(p_ah, font_name="Arial", size_pt=10, color_rgb=accent_rgb, bold=True)
            p_ah.space_after = Pt(6)
            
            fields_a = item_a.get("detail", "id (PK)").split(",")
            for field in fields_a[:5]:
                p = tf_a.add_paragraph()
                p.text = f"• {field.strip()}"
                apply_text_styling(p, font_name="Arial", size_pt=8.5, color_rgb=text_dark_rgb)
                
            item_b = next((it for it in items if it.get("category") == "entity_b"), items[1] if len(items) > 1 else (items[0] if items else {}))
            tbl_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.7) + x_shift, Inches(2.2), Inches(1.8), Inches(2.4))
            tbl_b.fill.solid()
            tbl_b.fill.fore_color.rgb = RGBColor(255, 255, 255)
            tbl_b.line.color.rgb = RGBColor(*accent_rgb)
            tbl_b.line.width = Pt(2)
            tf_b = tbl_b.text_frame
            tf_b.word_wrap = True
            p_bh = tf_b.paragraphs[0]
            p_bh.text = item_b.get("label", "Entity_B")
            apply_text_styling(p_bh, font_name="Arial", size_pt=10, color_rgb=accent_rgb, bold=True)
            p_bh.space_after = Pt(6)
            
            fields_b = item_b.get("detail", "id (PK)").split(",")
            for field in fields_b[:5]:
                p = tf_b.add_paragraph()
                p.text = f"• {field.strip()}"
                apply_text_styling(p, font_name="Arial", size_pt=8.5, color_rgb=text_dark_rgb)
                
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.6) + x_shift, Inches(3.4), Inches(1.1), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*accent_rgb)
            line.line.fill.background()
            return

        # 8. Comparison Matrix
        elif diagram_type == "matrix" and items:
            coords = [
                (Inches(7.8), Inches(1.8)),
                (Inches(10.2), Inches(1.8)),
                (Inches(7.8), Inches(4.1)),
                (Inches(10.2), Inches(4.1))
            ]
            for i, item in enumerate(items[:4]):
                x, y = coords[i]
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + x_shift, y, Inches(2.2), Inches(2.1))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
                card.line.color.rgb = RGBColor(*border_rgb)
                tf = card.text_frame
                tf.word_wrap = True
                
                p_h = tf.paragraphs[0]
                p_h.text = item.get("label", f"Quadrant {i+1}")
                apply_text_styling(p_h, font_name="Arial", size_pt=10, color_rgb=accent_rgb, bold=True)
                p_h.space_after = Pt(4)
                
                if item.get("detail"):
                    p_b = tf.add_paragraph()
                    p_b.text = trunc(item["detail"], 60)
                    apply_text_styling(p_b, font_name="Arial", size_pt=8.5, color_rgb=text_dark_rgb)
            return

    # --- Legacy Fallback logic ---
    # 1. Process Flow / Business Workflow
    if any(x in text for x in ["process", "workflow", "flowchart"]):
        top_offset = Inches(1.8)
        for i, point in enumerate(content[:4]):
            # Draw Step Rounded Rectangle
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3) + x_shift, top_offset, Inches(3.8), Inches(0.8))
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            shape.line.color.rgb = RGBColor(*border_rgb)
            tf = shape.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Step {i+1}: {trunc(point, 40)}"
            apply_text_styling(p, font_name="Arial", size_pt=11, color_rgb=text_dark_rgb, bold=True)
            
            top_offset += Inches(0.8)
            
            # Draw down connector arrow (unless last)
            if i < len(content[:4]) - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.1) + x_shift, top_offset, Inches(0.2), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(*accent_rgb)
                arrow.line.fill.background()
                top_offset += Inches(0.3)
        return

    # 2. Timeline
    if any(x in text for x in ["timeline", "schedule", "history"]):
        # Draw vertical timeline bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5) + x_shift, Inches(1.8), Inches(0.1), Inches(4.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*accent_rgb)
        bar.line.fill.background()
        
        top_offset = Inches(1.9)
        for i, point in enumerate(content[:4]):
            # Draw circle dot on line
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.35) + x_shift, top_offset + Inches(0.15), Inches(0.4), Inches(0.4))
            circle.fill.solid()
            circle.fill.fore_color.rgb = RGBColor(255, 255, 255)
            circle.line.color.rgb = RGBColor(*accent_rgb)
            circle.line.width = Pt(3)
            
            # Draw milestone card details
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9) + x_shift, top_offset, Inches(3.2), Inches(0.7))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            card.line.color.rgb = RGBColor(*border_rgb)
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = trunc(point, 45)
            apply_text_styling(p, font_name="Arial", size_pt=10, color_rgb=text_dark_rgb)
            
            top_offset += Inches(1.1)
        return

    # 3. Organizational Chart / Hierarchy
    if any(x in text for x in ["organ", "hierarchy", "chart"]):
        # Root Node
        root = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.2) + x_shift, Inches(1.8), Inches(2.0), Inches(0.7))
        root.fill.solid()
        root.fill.fore_color.rgb = RGBColor(*accent_rgb)
        root.line.fill.background()
        p_root = root.text_frame.paragraphs[0]
        p_root.text = "Root Level"
        apply_text_styling(p_root, font_name="Arial", size_pt=11, color_rgb=text_light_rgb, bold=True)
        
        # Vertical line down
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.15) + x_shift, Inches(2.5), Inches(0.05), Inches(0.6))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*accent_rgb)
        line.line.fill.background()
        
        # Horizontal branching line
        hbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3) + x_shift, Inches(3.1), Inches(3.8), Inches(0.05))
        hbar.fill.solid()
        hbar.fill.fore_color.rgb = RGBColor(*accent_rgb)
        hbar.line.fill.background()
        
        # Branch children boxes
        for i, point in enumerate(content[:3]):
            x_pos = Inches(7.8) + x_shift + i * Inches(1.4)
            vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_pos + Inches(0.55), Inches(3.15), Inches(0.05), Inches(0.45))
            vline.fill.solid()
            vline.fill.fore_color.rgb = RGBColor(*accent_rgb)
            vline.line.fill.background()
            
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, Inches(3.6), Inches(1.2), Inches(1.2))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            card.line.color.rgb = RGBColor(*border_rgb)
            tf = card.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"Div {i+1}:\n{trunc(point, 18)}"
            apply_text_styling(p, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
        return

    # 4. Mind Map
    if any(x in text for x in ["mind", "concept"]):
        # Central hub oval
        center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2) + x_shift, Inches(3.3), Inches(1.8), Inches(1.0))
        center.fill.solid()
        center.fill.fore_color.rgb = RGBColor(*accent_rgb)
        center.line.fill.background()
        p_c = center.text_frame.paragraphs[0]
        p_c.text = "Core Concept"
        apply_text_styling(p_c, font_name="Arial", size_pt=11, color_rgb=text_light_rgb, bold=True)
        
        spokes = [
            (Inches(7.8), Inches(1.8)),
            (Inches(10.6), Inches(1.8)),
            (Inches(9.2), Inches(4.8))
        ]
        for i, point in enumerate(content[:3]):
            x, y = spokes[i]
            # Draw line to spoke
            conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, (x + Inches(9.2))/2 + Inches(0.45) + x_shift, (y + Inches(3.3))/2 + Inches(0.25), Inches(0.05), Inches(0.4))
            conn.fill.solid()
            conn.fill.fore_color.rgb = RGBColor(*border_rgb)
            conn.line.fill.background()
            
            node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + x_shift, y, Inches(1.8), Inches(0.8))
            node.fill.solid()
            node.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            node.line.color.rgb = RGBColor(*border_rgb)
            tf = node.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = trunc(point, 25)
            apply_text_styling(p, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
        return

    # 5. Decision Tree
    if any(x in text for x in ["decision", "tree", "branch"]):
        # Decision diamond
        cond = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(9.2) + x_shift, Inches(1.8), Inches(2.0), Inches(1.0))
        cond.fill.solid()
        cond.fill.fore_color.rgb = RGBColor(*accent_rgb)
        cond.line.fill.background()
        p_cond = cond.text_frame.paragraphs[0]
        p_cond.text = "Decision"
        apply_text_styling(p_cond, font_name="Arial", size_pt=11, color_rgb=text_light_rgb, bold=True)
        
        # YES card (Green border)
        left_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8) + x_shift, Inches(3.8), Inches(1.8), Inches(1.2))
        left_card.fill.solid()
        left_card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
        left_card.line.color.rgb = RGBColor(34, 197, 94)
        left_card.line.width = Pt(2)
        tf_l = left_card.text_frame
        tf_l.word_wrap = True
        p_l = tf_l.paragraphs[0]
        p_l.text = f"[YES BRANCH]\n{trunc(content[0] if len(content) > 0 else 'Action A', 24)}"
        apply_text_styling(p_l, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
        
        # NO card (Red border)
        right_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6) + x_shift, Inches(3.8), Inches(1.8), Inches(1.2))
        right_card.fill.solid()
        right_card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
        right_card.line.color.rgb = RGBColor(239, 68, 68)
        right_card.line.width = Pt(2)
        tf_r = right_card.text_frame
        tf_r.word_wrap = True
        p_r = tf_r.paragraphs[0]
        p_r.text = f"[NO BRANCH]\n{trunc(content[1] if len(content) > 1 else 'Action B', 24)}"
        apply_text_styling(p_r, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
        return

    # 6. System Architecture (3-tier stack)
    if any(x in text for x in ["architecture", "infrastructure", "system"]):
        layers = [
            ("🖥️ Client UI", accent_light_rgb, text_dark_rgb),
            ("🚀 App API Server", accent_rgb, text_light_rgb),
            ("🗄️ Storage DB", accent_light_rgb, text_dark_rgb)
        ]
        
        top_offset = Inches(1.8)
        for i, (label, bg_rgb, fg_rgb) in enumerate(layers):
            layer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.3) + x_shift, top_offset, Inches(3.8), Inches(0.8))
            layer.fill.solid()
            layer.fill.fore_color.rgb = RGBColor(*bg_rgb)
            layer.line.color.rgb = RGBColor(*border_rgb)
            tf = layer.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            desc = content[i] if len(content) > i else "System components layer"
            p.text = f"{label}: {trunc(desc, 30)}"
            apply_text_styling(p, font_name="Arial", size_pt=10, color_rgb=fg_rgb, bold=True)
            
            top_offset += Inches(0.8)
            if i < 2:
                arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.1) + x_shift, top_offset, Inches(0.2), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = RGBColor(*accent_rgb)
                arrow.line.fill.background()
                top_offset += Inches(0.3)
        return

    # 7. Data Relationship / ERD
    if any(x in text for x in ["relationship", "database", "data", "erd"]):
        # Entity A
        tbl_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.8) + x_shift, Inches(2.2), Inches(1.8), Inches(2.2))
        tbl_a.fill.solid()
        tbl_a.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tbl_a.line.color.rgb = RGBColor(*accent_rgb)
        tbl_a.line.width = Pt(2)
        tf_a = tbl_a.text_frame
        tf_a.word_wrap = True
        p_ah = tf_a.paragraphs[0]
        p_ah.text = "Entity_A"
        apply_text_styling(p_ah, font_name="Arial", size_pt=11, color_rgb=accent_rgb, bold=True)
        p_ah.space_after = Pt(6)
        for field in ["id (PK)", "name", "desc"]:
            p = tf_a.add_paragraph()
            p.text = f"• {field}"
            apply_text_styling(p, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
            
        # Entity B
        tbl_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.7) + x_shift, Inches(2.2), Inches(1.8), Inches(2.2))
        tbl_b.fill.solid()
        tbl_b.fill.fore_color.rgb = RGBColor(255, 255, 255)
        tbl_b.line.color.rgb = RGBColor(*accent_rgb)
        tbl_b.line.width = Pt(2)
        tf_b = tbl_b.text_frame
        tf_b.word_wrap = True
        p_bh = tf_b.paragraphs[0]
        p_bh.text = "Entity_B"
        apply_text_styling(p_bh, font_name="Arial", size_pt=11, color_rgb=accent_rgb, bold=True)
        p_bh.space_after = Pt(6)
        for field in ["id (PK)", "fk_id (FK)", "value"]:
            p = tf_b.add_paragraph()
            p.text = f"• {field}"
            apply_text_styling(p, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)

        # Connector relation line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.6) + x_shift, Inches(3.2), Inches(1.1), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*accent_rgb)
        line.line.fill.background()
        return

    # 8. Comparison Matrix (2x2 grid)
    if any(x in text for x in ["matrix", "compare", "grid", "quadrant", "table"]):
        coords = [
            (Inches(7.8), Inches(1.8)),
            (Inches(10.2), Inches(1.8)),
            (Inches(7.8), Inches(4.1)),
            (Inches(10.2), Inches(4.1))
        ]
        for i, point in enumerate(content[:4]):
            x, y = coords[i]
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + x_shift, y, Inches(2.2), Inches(2.0))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*accent_light_rgb)
            card.line.color.rgb = RGBColor(*border_rgb)
            tf = card.text_frame
            tf.word_wrap = True
            
            p_h = tf.paragraphs[0]
            p_h.text = f"Quadrant {i+1}"
            apply_text_styling(p_h, font_name="Arial", size_pt=11, color_rgb=accent_rgb, bold=True)
            p_h.space_after = Pt(4)
            
            p_b = tf.add_paragraph()
            p_b.text = trunc(point, 50)
            apply_text_styling(p_b, font_name="Arial", size_pt=9, color_rgb=text_dark_rgb)
        return

    # Default fallback: visual layout suggestion textbox
    visual_box = slide.shapes.add_textbox(Inches(7.8) + x_shift, Inches(1.8), Inches(4.7), Inches(4.8))
    tf_visual = visual_box.text_frame
    tf_visual.word_wrap = True
    
    p_vis_lbl = tf_visual.paragraphs[0]
    p_vis_lbl.text = "RECOMMENDED VISUAL LAYOUT:"
    apply_text_styling(p_vis_lbl, font_name="Arial", size_pt=12, color_rgb=accent_rgb, bold=True)
    p_vis_lbl.space_after = Pt(10)
    
    p_vis_val = tf_visual.add_paragraph()
    p_vis_val.text = suggestion
    apply_text_styling(p_vis_val, font_name="Arial", size_pt=14, color_rgb=(71, 85, 105), italic=True)

def main():
    if len(sys.argv) < 3:
        print("Usage: python3 ppt_exporter.py <slides_json_or_file> <output_pptx_path>")
        sys.exit(1)

    json_input = sys.argv[1]
    output_path = sys.argv[2]

    # Load JSON from file or parse directly
    try:
        if os.path.exists(json_input):
            with open(json_input, 'r', encoding='utf-8') as f:
                slides = json.load(f)
        else:
            slides = json.loads(json_input)
    except Exception as e:
        print(f"Error loading JSON input: {e}", file=sys.stderr)
        sys.exit(1)

    # Allow slides key wrapping or direct array
    if isinstance(slides, dict) and "slides" in slides:
        slides = slides["slides"]

    theme_name = sys.argv[3] if len(sys.argv) > 3 else "slate"
    palette = THEME_PALETTES.get(theme_name, THEME_PALETTES["slate"])

    prs = Presentation()
    # Set slide size to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout for custom placement to make it look premium
    blank_layout = prs.slide_layouts[6]

    for idx, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        
        title_text = slide_data.get("title", f"Slide {idx + 1}")
        content_points = slide_data.get("content", [])
        visual_suggestion = slide_data.get("visual_suggestion", "")

        if idx == 0:
            # Title Slide
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*palette["bg"])
            
            # Extract image URL and layout position from slide data
            image_data = slide_data.get("image")
            image_url = None
            image_pos = "none"
            if image_data and isinstance(image_data, dict):
                image_url = image_data.get("url")
                image_pos = image_data.get("position", "none")
                if not image_url:
                    image_pos = "none"

            # Separate Title and Subtitle text boxes to prevent overlaps and overflows
            title_size = 36
            subtitle_size = 14
            
            if image_pos != "none":
                if image_pos == "left":
                    add_image_from_url(slide, image_url, Inches(0.8), Inches(1.5), Inches(4.7), Inches(5.1))
                    txBox_title = slide.shapes.add_textbox(Inches(6.0), Inches(1.6), Inches(6.5), Inches(2.5))
                    txBox_sub = slide.shapes.add_textbox(Inches(6.0), Inches(4.4), Inches(6.5), Inches(2.2))
                elif image_pos == "top":
                    add_image_from_url(slide, image_url, Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.2))
                    txBox_title = slide.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(11.333), Inches(1.4))
                    txBox_sub = slide.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(11.333), Inches(1.2))
                    title_size = 32
                else:  # "right"
                    add_image_from_url(slide, image_url, Inches(7.8), Inches(1.5), Inches(4.7), Inches(5.1))
                    txBox_title = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(6.2), Inches(2.5))
                    txBox_sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(6.2), Inches(2.2))
            else:
                # Text box for title & description (text-only title slide)
                txBox_title = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.2))
                txBox_sub = slide.shapes.add_textbox(Inches(1.0), Inches(4.4), Inches(11.333), Inches(2.2))
                title_size = 40
                subtitle_size = 16
                
            # Render Title
            tf_title = txBox_title.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            apply_text_styling(p_title, font_name="Georgia", size_pt=title_size, color_rgb=palette["text"], bold=True)
            
            # Render Subtitle (Content points)
            if content_points:
                tf_sub = txBox_sub.text_frame
                tf_sub.word_wrap = True
                p_sub = tf_sub.paragraphs[0]
                p_sub.text = " • ".join(content_points)
                apply_text_styling(p_sub, font_name="Arial", size_pt=subtitle_size, color_rgb=palette["muted"])
        else:
            # Content Slide
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*palette["bg"])
            
            # Slide Header Title (Font size reduced from 32 to 24 to prevent wraps and title overlaps)
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
            tf_title = title_box.text_frame
            tf_title.word_wrap = True
            p_title = tf_title.paragraphs[0]
            p_title.text = title_text
            apply_text_styling(p_title, font_name="Georgia", size_pt=24, color_rgb=palette["text"], bold=True)
            
            # Extract image URL and layout position from slide data
            image_data = slide_data.get("image")
            image_url = None
            image_pos = "none"
            if image_data and isinstance(image_data, dict):
                image_url = image_data.get("url")
                image_pos = image_data.get("position", "none")
                if not image_url:
                    image_pos = "none"

            # Determine alternating columns for PPT Designer layout
            is_even_slide = (idx % 2 == 0)
            content_left = Inches(6.0) if (is_even_slide and visual_suggestion) else Inches(0.8)
            
            # 1. Slide has BOTH Diagram and Image
            if visual_suggestion and image_pos != "none":
                # Stacking: Text is at top-half of text column, Image is at bottom-half of text column
                # Shifted content top to 1.8, and restricted height to 2.7 to avoid image overlaps
                content_box = slide.shapes.add_textbox(content_left, Inches(1.8), Inches(6.5), Inches(2.7))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True
                for p_idx, point in enumerate(content_points):
                    p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                    p.text = f"•  {point}"
                    # Reduced bullet size to 12pt and spacing to Pt(3) to prevent overflow
                    apply_text_styling(p, font_name="Arial", size_pt=12, color_rgb=palette["text"])
                    p.space_after = Pt(3)
                
                # Image placed further down at top=4.7 with height=2.0 to avoid overlap
                add_image_from_url(slide, image_url, content_left, Inches(4.7), Inches(6.5), Inches(2.0))
                
                # Diagram card shifted to top=1.8 with height=4.9
                x_shift = -Inches(7.0) if is_even_slide else Inches(0)
                card_left = Inches(7.8) + x_shift
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, Inches(1.8), Inches(4.7), Inches(4.9))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*palette["cardBg"])
                card.line.color.rgb = RGBColor(*palette["border"])
                card.line.width = Pt(1.5)
                
                draw_diagram(slide, visual_suggestion, content_points, slide_data.get("diagram"), theme_name, x_shift)
                
            # 2. Slide has Diagram but NO Image
            elif visual_suggestion:
                content_box = slide.shapes.add_textbox(content_left, Inches(1.8), Inches(6.5), Inches(4.9))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True
                for p_idx, point in enumerate(content_points):
                    p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                    p.text = f"•  {point}"
                    apply_text_styling(p, font_name="Arial", size_pt=16, color_rgb=palette["text"])
                    p.space_after = Pt(14)
                
                x_shift = -Inches(7.0) if is_even_slide else Inches(0)
                card_left = Inches(7.8) + x_shift
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, card_left, Inches(1.8), Inches(4.7), Inches(4.9))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*palette["cardBg"])
                card.line.color.rgb = RGBColor(*palette["border"])
                card.line.width = Pt(1.5)
                
                draw_diagram(slide, visual_suggestion, content_points, slide_data.get("diagram"), theme_name, x_shift)
                
            # 3. Slide has NO Diagram but HAS Image
            elif image_pos != "none":
                if image_pos == "left":
                    # Render image in left column, text in right column
                    add_image_from_url(slide, image_url, Inches(0.8), Inches(1.8), Inches(4.7), Inches(4.9))
                    
                    content_box = slide.shapes.add_textbox(Inches(6.0), Inches(1.8), Inches(6.5), Inches(4.9))
                    tf_content = content_box.text_frame
                    tf_content.word_wrap = True
                    for p_idx, point in enumerate(content_points):
                        p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                        p.text = f"•  {point}"
                        apply_text_styling(p, font_name="Arial", size_pt=16, color_rgb=palette["text"])
                        p.space_after = Pt(14)
                elif image_pos == "top":
                    # Render image on top, text below
                    add_image_from_url(slide, image_url, Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.2))
                    
                    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.733), Inches(2.5))
                    tf_content = content_box.text_frame
                    tf_content.word_wrap = True
                    for p_idx, point in enumerate(content_points):
                        p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                        p.text = f"•  {point}"
                        apply_text_styling(p, font_name="Arial", size_pt=16, color_rgb=palette["text"])
                        p.space_after = Pt(14)
                else:  # default or "right"
                    # Render text in left column, image in right column
                    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6.5), Inches(4.9))
                    tf_content = content_box.text_frame
                    tf_content.word_wrap = True
                    for p_idx, point in enumerate(content_points):
                        p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                        p.text = f"•  {point}"
                        apply_text_styling(p, font_name="Arial", size_pt=16, color_rgb=palette["text"])
                        p.space_after = Pt(14)
                    
                    add_image_from_url(slide, image_url, Inches(7.8), Inches(1.8), Inches(4.7), Inches(4.9))
                
            # 4. Slide has NEITHER Diagram nor Image
            else:
                # Full width text block
                content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.9))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True
                for p_idx, point in enumerate(content_points):
                    p = tf_content.add_paragraph() if p_idx > 0 else tf_content.paragraphs[0]
                    p.text = f"•  {point}"
                    apply_text_styling(p, font_name="Arial", size_pt=18, color_rgb=palette["text"])
                    p.space_after = Pt(16)

    prs.save(output_path)
    print(f"PowerPoint saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
